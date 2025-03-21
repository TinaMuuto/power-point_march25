import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import re
import requests
from PIL import Image
from copy import deepcopy
import math
from concurrent.futures import ThreadPoolExecutor, as_completed

# Filstier – tilpas efter behov
MAPPING_FILE_PATH = "mapping-file.xlsx"
STOCK_FILE_PATH = "stock.xlsx"
TEMPLATE_FILE_PATH = "template-generator.pptx"

# --- Forventede kolonner i mapping-fil ---
REQUIRED_MAPPING_COLS_ORIG = [
    "{{Product name}}",
    "{{Product code}}",
    "{{Product country of origin}}",
    "{{Product height}}",
    "{{Product width}}",
    "{{Product length}}",
    "{{Product depth}}",
    "{{Product seat height}}",
    "{{Product diameter}}",
    "{{CertificateName}}",
    "{{Product Consumption COM}}",
    "{{Product Fact Sheet link}}",
    "{{Product configurator link}}",
    "{{Product Packshot1}}",
    "{{Product Lifestyle1}}",
    "{{Product Lifestyle2}}",
    "{{Product Lifestyle3}}",
    "{{Product Lifestyle4}}",
    "ProductKey"
]

# --- Forventede kolonner i stock-fil ---
REQUIRED_STOCK_COLS_ORIG = [
    "productkey",
    "variantname",
    "rts",
    "mto"
]

# --- Placeholders til erstatning i templaten ---
TEXT_PLACEHOLDERS_ORIG = {
    "{{Product name}}": "Product Name:",
    "{{Product code}}": "Product Code:",
    "{{Product country of origin}}": "Country of origin:",
    "{{Product height}}": "Height:",
    "{{Product width}}": "Width:",
    "{{Product length}}": "Length:",
    "{{Product depth}}": "Depth:",
    "{{Product seat height}}": "Seat Height:",
    "{{Product diameter}}": "Diameter:",
    "{{CertificateName}}": "Test & certificates for the product:",
    "{{Product Consumption COM}}": "Consumption information for COM:"
}

HYPERLINK_PLACEHOLDERS_ORIG = {
    "{{Product Fact Sheet link}}": "Download Product Fact Sheet",
    "{{Product configurator link}}": "Click to configure product"
}

IMAGE_PLACEHOLDERS_ORIG = [
    "{{Product Packshot1}}",
    "{{Product Lifestyle1}}",
    "{{Product Lifestyle2}}",
    "{{Product Lifestyle3}}",
    "{{Product Lifestyle4}}",
]

# --- Almindelig gruppering af variantnavne (simpel version) ---
def group_by_color_and_size(variant_names):
    """
    Denne funktion grupperer variantnavnene efter farve og samler unikke størrelser.
    For hvert variantnavn forventes formatet: "Farve - [noget] - Størrelse".
    Den bruger den første del som farve og den sidste som størrelse.
    Output: "Farve: Størrelse 1, Størrelse 2, ..."
    Hvis der ikke findes en " - " separator, returneres navnet uændret.
    """
    groups = {}
    for name in variant_names:
        if " - " in name:
            parts = name.split(" - ")
            color = parts[0].strip()
            size = parts[-1].strip()  # den sidste del
        else:
            color = name.strip()
            size = ""
        groups.setdefault(color, set())
        if size:
            groups[color].add(size)
    output_lines = []
    for color, sizes in groups.items():
        if sizes:
            output_lines.append(f"{color}: {', '.join(sorted(sizes))}")
        else:
            output_lines.append(color)
    return "\n".join(sorted(output_lines))

# --- Alternativ logik via en fast defineret produktkonfigurator ---
class ProductConfigurator:
    def __init__(self):
        # Faste produktkombinationer: (overflade, mellemstykke, ben) -> benfarver
        self.products = {
            ("Black Linoleum", "Plywood", "Plywood"): ["Black", "Grey", "Sand", "White"],
            ("Grey Linoleum", "Plywood", "Plywood"): ["Black", "Grey", "Sand", "White"],
            ("Oak Lacquered Oak Veneer", "Plywood", "Plywood"): ["Black", "Grey", "Sand", "White"],
            ("Oak Oiled Oak", "Oak Oiled Oak", "Oak Oiled Oak"): ["Black", "Grey", "Sand", "White"],
            ("Sand Laminate", "Plywood", "Plywood"): ["Black", "Grey", "Sand", "White"],
            ("Smoked Oak Oiled Oak", "Smoked Oak Oiled Oak", "Smoked Oak Oiled Oak"): ["Black", "Grey", "Sand", "White"],
            ("White Laminate", "Plywood", "Plywood"): ["Black", "Grey", "Sand", "White"],
        }
        # Standardstørrelser
        self.default_sizes = [
            "170 x 85 cm / 67 x 33.5\"",
            "225 x 90 cm / 88.5 x 35.5\"",
            "255 x 108 cm / 100.5 x 42.5\"",
            "295 x 108 cm / 116 x 42.5\"",
        ]
        # Specifikke størrelser for enkelte overflader
        self.specific_sizes = {
            "Sand Laminate": [
                "225 x 90 cm / 88.5 x 35.5\"",
                "255 x 108 cm / 100.5 x 42.5\"",
                "295 x 108 cm / 116 x 42.5\"",
            ]
        }
    
    def get_options(self, product_name):
        """Hvis produktnavnet indeholder et af overfladenavnene, returnér de fast definerede benfarver og størrelser."""
        for (surface, core, legs), colors in self.products.items():
            if surface.lower() in product_name.lower():
                sizes = self.specific_sizes.get(surface, self.default_sizes)
                return {"benfarver": colors, "størrelser": sizes}
        return None

configurator = ProductConfigurator()

# --- Almindelige hjælpefunktioner ---
def normalize_text(s):
    return re.sub(r"\s+", "", str(s).replace("\u00A0", " ")).lower()

def normalize_col(col):
    return normalize_text(col)

def find_mapping_row(item_no, mapping_df, mapping_prod_key):
    norm_item = normalize_text(item_no)
    for idx, row in mapping_df.iterrows():
        code = row.get(mapping_prod_key, "")
        if normalize_text(code) == norm_item:
            return row
    if "-" in str(item_no):
        partial = normalize_text(item_no.split("-")[0])
        for idx, row in mapping_df.iterrows():
            code = row.get(mapping_prod_key, "")
            if normalize_text(code).startswith(partial):
                return row
    return None

def process_stock_rts_alternative(mapping_row, stock_df):
    product_key = mapping_row.get("productkey", "")
    if not product_key or pd.isna(product_key):
        return ""
    norm_product_key = normalize_text(product_key)
    try:
        filtered = stock_df[stock_df["productkey"].apply(lambda x: normalize_text(x) == norm_product_key)]
    except KeyError as e:
        st.error(f"KeyError i RTS: {e}")
        return ""
    if filtered.empty:
        return ""
    filtered = filtered[filtered["rts"].notna() & (filtered["rts"] != "")]
    if filtered.empty:
        return ""
    try:
        variant_names = filtered["variantname"].dropna().astype(str).tolist()
    except KeyError as e:
        st.error(f"KeyError i RTS variantname: {e}")
        return ""
    unique_variant_names = list(dict.fromkeys(variant_names))
    
    # Tjek, om produktnavnet matcher en af de konfigurationer
    product_name = str(mapping_row.get(normalize_col("{{Product name}}"), ""))
    options = configurator.get_options(product_name)
    if options:
        return f"Benfarver: {', '.join(options['benfarver'])}\nStørrelser: {', '.join(options['størrelser'])}"
    else:
        return group_by_color_and_size(unique_variant_names)

def process_stock_mto_alternative(mapping_row, stock_df):
    product_key = mapping_row.get("productkey", "")
    if not product_key or pd.isna(product_key):
        return ""
    norm_product_key = normalize_text(product_key)
    try:
        filtered = stock_df[stock_df["productkey"].apply(lambda x: normalize_text(x) == norm_product_key)]
    except KeyError as e:
        st.error(f"KeyError i MTO: {e}")
        return ""
    if filtered.empty:
        return ""
    filtered = filtered[filtered["mto"].notna() & (filtered["mto"] != "")]
    if filtered.empty:
        return ""
    try:
        variant_names = filtered["variantname"].dropna().astype(str).tolist()
    except KeyError as e:
        st.error(f"KeyError i MTO variantname: {e}")
        return ""
    unique_variant_names = list(dict.fromkeys(variant_names))
    
    product_name = str(mapping_row.get(normalize_col("{{Product name}}"), ""))
    options = configurator.get_options(product_name)
    if options:
        return f"Benfarver: {', '.join(options['benfarver'])}\nStørrelser: {', '.join(options['størrelser'])}"
    else:
        return group_by_color_and_size(unique_variant_names)

@st.cache_data(show_spinner=False)
def fetch_and_process_image_cached(url, quality=70, max_size=(1200, 1200)):
    try:
        response = requests.get(url, timeout=30)
        if response.status_code == 200:
            img = Image.open(io.BytesIO(response.content))
            if img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info) or (img.format and img.format.lower() == "tiff"):
                img = img.convert("RGB")
            img.thumbnail(max_size, Image.LANCZOS)
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format="JPEG", quality=quality, optimize=True)
            img_byte_arr.seek(0)
            return img_byte_arr
    except Exception as e:
        st.warning(f"Fejl ved hentning af billede fra {url}: {e}")
    return None

def replace_image_placeholders_parallel(slide, image_values):
    from concurrent.futures import ThreadPoolExecutor, as_completed
    tasks = []
    for ph in IMAGE_PLACEHOLDERS_ORIG:
        norm_ph = normalize_text(ph)
        url = image_values.get(ph, "")
        if url:
            tasks.append((ph, url))
    results = {}
    with ThreadPoolExecutor(max_workers=5) as executor:
        future_to_ph = {executor.submit(fetch_and_process_image_cached, url): ph for ph, url in tasks}
        for future in as_completed(future_to_ph):
            ph = future_to_ph[future]
            try:
                results[ph] = future.result()
            except Exception as exc:
                st.warning(f"Fejl ved parallel billedhentning for {ph}: {exc}")
                results[ph] = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            tekst = shape.text
            for ph in IMAGE_PLACEHOLDERS_ORIG:
                norm_ph = normalize_text(ph)
                if norm_ph in normalize_text(tekst):
                    url = image_values.get(ph, "")
                    if url and results.get(ph):
                        img = Image.open(results[ph])
                        original_width, original_height = img.size
                        target_width = shape.width
                        target_height = shape.height
                        scale = min(target_width / original_width, target_height / original_height)
                        new_width = int(original_width * scale)
                        new_height = int(original_height * scale)
                        new_img_stream = io.BytesIO()
                        img.save(new_img_stream, format="JPEG")
                        new_img_stream.seek(0)
                        slide.shapes.add_picture(new_img_stream, shape.left, shape.top, width=new_width, height=new_height)
                        shape.text = ""
                    break

def replace_hyperlink_placeholders(slide, hyperlink_values):
    import re
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, (display_text, url) in hyperlink_values.items():
                        key = placeholder.strip("{}").strip()
                        pattern = r"\{\{\s*" + re.escape(key) + r"\s*\}\}"
                        if re.search(pattern, run.text):
                            run.text = re.sub(pattern, display_text, run.text)
                            try:
                                run.hyperlink.address = url
                            except Exception as e:
                                st.warning(f"Hyperlink for {placeholder} kunne ikke indsættes: {e}")

def replace_text_placeholders(slide, placeholder_values):
    import re
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join([run.text for run in paragraph.runs])
                new_text = full_text
                for placeholder, replacement in placeholder_values.items():
                    key = placeholder.strip("{}").strip()
                    pattern = r"\{\{\s*" + re.escape(key) + r"\s*\}\}"
                    new_text = re.sub(pattern, replacement, new_text)
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    for i in range(len(paragraph.runs)-1, -1, -1):
                        paragraph.runs[i].text = ""
                    first_run.text = new_text

def duplicate_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    new_slide.shapes._spTree.clear()
    for shape in slide.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape._element))
    # Fjern eventuelle hidden-tags, så sliden vises korrekt
    for elem in new_slide._element.xpath('.//p:hiddenslide'):
        elem.getparent().remove(elem)
    return new_slide

def delete_slide(prs, slide_index):
    slide_id = prs.slides._sldIdLst[slide_index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(slide_id)

# --- Main Streamlit App ---
def main():
    st.title("PowerPoint Generator App")
    st.write("Indsæt varenumre (Item no) – ét pr. linje:")
    st.info("Bemærk: Indsæt varenumre uden ekstra mellemrum omkring bindestreger, f.eks. '03194', '03094', osv.")
    
    pasted_text = st.text_area("Indsæt varenumre her", height=200)
    if not pasted_text.strip():
        st.error("Indsæt venligst varenumre i tekstfeltet.")
        return

    varenumre = [line.strip() for line in pasted_text.splitlines() if line.strip()]
    if not varenumre:
        st.error("Ingen gyldige varenumre fundet.")
        return

    user_df = pd.DataFrame({"Item no": varenumre, "Product name": [""] * len(varenumre)})
    
    progress_bar = st.progress(0)
    status = st.empty()
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Filer uploadet og brugerdata oprettet.</div>", unsafe_allow_html=True)
    progress_bar.progress(10)
    
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Indlæser mapping-fil...</div>", unsafe_allow_html=True)
    try:
        mapping_df = pd.read_excel(MAPPING_FILE_PATH)
        mapping_df.columns = [normalize_col(col) for col in mapping_df.columns]
    except Exception as e:
        st.error(f"Fejl ved læsning af mapping-fil: {e}")
        return
    normalized_required_mapping_cols = [normalize_col(col) for col in REQUIRED_MAPPING_COLS_ORIG]
    missing_mapping_cols = [req for req in normalized_required_mapping_cols if req not in mapping_df.columns]
    if missing_mapping_cols:
        st.error(f"Mapping-filen mangler kolonner: {missing_mapping_cols}.")
        return
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Mapping-fil indlæst.</div>", unsafe_allow_html=True)
    progress_bar.progress(30)
    MAPPING_PRODUCT_CODE_KEY = normalize_col("{{Product code}}")
    
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Indlæser stock-fil...</div>", unsafe_allow_html=True)
    try:
        stock_df = pd.read_excel(STOCK_FILE_PATH)
        stock_df.columns = [normalize_col(col) for col in stock_df.columns]
    except Exception as e:
        st.error(f"Fejl ved læsning af stock-fil: {e}")
        return
    normalized_required_stock_cols = [normalize_col(col) for col in REQUIRED_STOCK_COLS_ORIG]
    missing_stock_cols = [req for req in normalized_required_stock_cols if req not in stock_df.columns]
    if missing_stock_cols:
        st.error(f"Stock-filen mangler kolonner: {missing_stock_cols}.")
        return
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Stock-fil indlæst.</div>", unsafe_allow_html=True)
    progress_bar.progress(50)
    
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Indlæser PowerPoint-template...</div>", unsafe_allow_html=True)
    try:
        prs = Presentation(TEMPLATE_FILE_PATH)
    except Exception as e:
        st.error(f"Fejl ved læsning af template-fil: {e}")
        return
    if len(prs.slides) < 1:
        st.error("Template-filen skal indeholde mindst én slide.")
        return
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Template-fil indlæst.</div>", unsafe_allow_html=True)
    progress_bar.progress(70)
    
    # Lav en kopi af den originale templateslide og slet den originale
    template_slide = prs.slides[0]
    template_copy = deepcopy(template_slide)
    delete_slide(prs, 0)
    
    total_products = len(user_df)
    batch_size = 10
    num_batches = math.ceil(total_products / batch_size)
    status.markdown(f"<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: {total_products} varer opdelt i {num_batches} batch(es).</div>", unsafe_allow_html=True)
    
    missing_items = []
    for batch_index in range(num_batches):
        status.markdown(f"<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Behandler batch {batch_index + 1} af {num_batches}...</div>", unsafe_allow_html=True)
        batch_df = user_df.iloc[batch_index * batch_size : (batch_index + 1) * batch_size]
        for idx, product in batch_df.iterrows():
            item_no = product["Item no"]
            slide = duplicate_slide(prs, template_copy)
            mapping_row = find_mapping_row(item_no, mapping_df, MAPPING_PRODUCT_CODE_KEY)
            if mapping_row is None:
                missing_items.append(item_no)
                placeholder_texts = {}
                for ph, label in TEXT_PLACEHOLDERS_ORIG.items():
                    if ph == "{{Product code}}":
                        placeholder_texts[ph] = f"{label} {item_no}"
                    else:
                        placeholder_texts[ph] = ""
                placeholder_texts["{{Product RTS}}"] = "Product in stock versions:\n\n"
                placeholder_texts["{{Product MTO}}"] = "Avilable for made to order:\n\n"
                replace_text_placeholders(slide, placeholder_texts)
                replace_hyperlink_placeholders(slide, {})
            else:
                placeholder_texts = {}
                for ph, label in TEXT_PLACEHOLDERS_ORIG.items():
                    norm_ph = normalize_col(ph)
                    value = mapping_row.get(norm_ph, "")
                    if pd.isna(value) or not str(value).strip():
                        placeholder_texts[ph] = ""
                    else:
                        if ph in ("{{Product code}}", "{{Product name}}", "{{Product country of origin}}"):
                            placeholder_texts[ph] = f"{label} {value}"
                        elif ph in ("{{CertificateName}}", "{{Product Consumption COM}}"):
                            placeholder_texts[ph] = f"{label}\n\n{value}"
                        else:
                            placeholder_texts[ph] = f"{label}\n{value}"
                
                rts_text = process_stock_rts_alternative(mapping_row, stock_df)
                mto_text = process_stock_mto_alternative(mapping_row, stock_df)
                placeholder_texts["{{Product RTS}}"] = f"Product in stock versions:\n\n{rts_text}"
                placeholder_texts["{{Product MTO}}"] = f"Avilable for made to order:\n\n{mto_text}"
                
                replace_text_placeholders(slide, placeholder_texts)
                
                hyperlink_vals = {}
                for ph, display_text in HYPERLINK_PLACEHOLDERS_ORIG.items():
                    norm_ph = normalize_col(ph)
                    url = mapping_row.get(norm_ph, "")
                    if pd.isna(url) or not str(url).strip():
                        url = ""
                    hyperlink_vals[ph] = (display_text, url)
                replace_hyperlink_placeholders(slide, hyperlink_vals)
                
                image_vals = {}
                for ph in IMAGE_PLACEHOLDERS_ORIG:
                    norm_ph = normalize_col(ph)
                    url = mapping_row.get(norm_ph, "")
                    if pd.isna(url) or not str(url).strip():
                        url = ""
                    image_vals[ph] = url
                replace_image_placeholders_parallel(slide, image_vals)
        progress = 70 + int((batch_index + 1) / num_batches * 30)
        progress_bar.progress(progress)
    
    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: Generering fuldført!</div>", unsafe_allow_html=True)
    ppt_io = io.BytesIO()
    try:
        prs.save(ppt_io)
        ppt_io.seek(0)
    except Exception as e:
        st.error(f"Fejl ved gemning af PowerPoint: {e}")
        return

    status.markdown("<div style='background-color:#f0f0f0; padding: 10px; border-radius: 5px;'>Status: PowerPoint genereret succesfuldt!</div>", unsafe_allow_html=True)
    st.success("PowerPoint genereret succesfuldt!")
    st.download_button("Download PowerPoint", ppt_io,
                       file_name="generated_presentation.pptx",
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    
    if missing_items:
        st.text_area("Manglende varenumre (kopier her):", value="\n".join(missing_items), height=100)
    
    st.session_state.generated_ppt = ppt_io

if __name__ == '__main__':
    if 'generated_ppt' not in st.session_state:
        st.session_state.generated_ppt = None
    main()
